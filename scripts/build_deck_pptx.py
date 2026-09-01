#!/usr/bin/env python3
"""Build the Agentic Graph Analytics intro deck as an editable .pptx.

Mirrors the content of ``docs/deck-agentic-graph-analytics.html`` in a format
that can be presented and edited in PowerPoint, Keynote or Google Slides.

Written as a generator rather than a hand-built file so the deck can be
regenerated when the content changes:

    python scripts/build_deck_pptx.py

Fonts are deliberately limited to faces present on both macOS and Windows —
a .pptx that silently falls back to Calibri on the reviewer's machine looks
broken in a way that is hard to notice before you are on stage.
"""

from __future__ import annotations

import functools
import pathlib
import textwrap

from PIL import ImageFont

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = pathlib.Path(__file__).resolve().parent.parent / "docs" / "agentic-graph-analytics.pptx"

# ---------------------------------------------------------------- palette
PAPER = RGBColor(0xF7, 0xF8, 0xF6)
RAISED = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x11, 0x16, 0x1A)
INK_SOFT = RGBColor(0x3D, 0x47, 0x4D)
INK_FAINT = RGBColor(0x73, 0x7E, 0x84)
RULE = RGBColor(0xDB, 0xD9, 0xD2)
ACCENT = RGBColor(0x0F, 0x6E, 0x5C)
ACCENT_WASH = RGBColor(0xE4, 0xEF, 0xEB)
CAUTION = RGBColor(0x96, 0x60, 0x14)
CAUTION_WASH = RGBColor(0xF6, 0xEE, 0xDF)

DISPLAY = "Helvetica Neue"   # falls back to Helvetica / Arial
BODY = "Georgia"
MONO = "Menlo"               # falls back to Consolas / Courier New

# ---------------------------------------------------------------- geometry
W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.85)
CONTENT_W = W - 2 * MARGIN
FOOT_Y = H - Inches(0.62)


# ---------------------------------------------------------------- measurement
# Earlier revisions guessed line counts from character length. The guesses were
# wrong in both directions, so blocks either collided with the next one or left
# a gap. These measure with the real font files instead, which is why the
# vertical rhythm below can be computed rather than hand-tuned per slide.
_FONT_FILES = {
    DISPLAY: "/System/Library/Fonts/Helvetica.ttc",
    BODY: "/System/Library/Fonts/Supplemental/Georgia.ttf",
    MONO: "/System/Library/Fonts/Supplemental/Courier New.ttf",
}
_PX_PER_PT = 4.0  # measure at 4x for sub-point accuracy


@functools.lru_cache(maxsize=256)
def _pil_font(family: str, size_pt: float, bold: bool):
    path = _FONT_FILES.get(family, _FONT_FILES[BODY])
    try:
        if family == DISPLAY:
            # Helvetica.ttc: index 0 regular, 1 bold.
            return ImageFont.truetype(path, int(size_pt * _PX_PER_PT), index=1 if bold else 0)
        return ImageFont.truetype(path, int(size_pt * _PX_PER_PT))
    except OSError:
        return ImageFont.load_default()


def measure_lines(text: str, *, font=BODY, size=14, bold=False, width_in=10.0) -> int:
    """Number of wrapped lines `text` needs at `size` inside `width_in` inches."""
    pil = _pil_font(font, size, bold)
    limit = width_in * 72 * _PX_PER_PT
    total = 0
    for hard_line in text.split("\n"):
        words, line, count = hard_line.split(), "", 1
        for word in words:
            trial = f"{line} {word}".strip()
            if pil.getlength(trial) <= limit or not line:
                line = trial
            else:
                count += 1
                line = word
        total += count
    return max(1, total)


def block_height(text, *, font=BODY, size=14, bold=False, width_in=10.0, line_spacing=1.35):
    """Rendered height of a text block, in EMU."""
    return Inches(measure_lines(text, font=font, size=size, bold=bold,
                                width_in=width_in) * size * line_spacing / 72)


def _text(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(paragraph, text, *, font=BODY, size=14, bold=False, italic=False,
         color=INK_SOFT, spacing=None):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if spacing is not None:
        # python-pptx has no letter-spacing API; set it on the run's rPr.
        run.font._element.set("spc", str(int(spacing * 100)))
    return run


def _para(tf, *, first=False, space_after=0, space_before=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    return p


def new_slide(prs, *, number, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER

    # running foot: hairline + slide position + section name
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, FOOT_Y, CONTENT_W, Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False

    tf = _text(slide, MARGIN, FOOT_Y + Inches(0.10), Inches(3), Inches(0.3))
    _run(_para(tf, first=True), f"{number:02d} / 11", font=MONO, size=9,
         color=INK_FAINT, spacing=1.2)

    tf = _text(slide, W - MARGIN - Inches(4), FOOT_Y + Inches(0.10), Inches(4), Inches(0.3))
    p = _para(tf, first=True)
    p.alignment = PP_ALIGN.RIGHT
    _run(p, section.upper(), font=MONO, size=9, color=INK_FAINT, spacing=1.2)
    return slide


def eyebrow(slide, text, y=Inches(0.95), color=INK_FAINT):
    tf = _text(slide, MARGIN, y, CONTENT_W, Inches(0.3))
    _run(_para(tf, first=True), text.upper(), font=MONO, size=10,
         color=color, spacing=1.6)
    return y + Inches(0.42)


def headline(slide, text, y, *, size=38, width=None, color=INK, gap=0.30):
    w = width or CONTENT_W
    h = block_height(text, font=DISPLAY, size=size, bold=True,
                     width_in=w / Inches(1), line_spacing=1.06)
    tf = _text(slide, MARGIN, y, w, h + Inches(0.1))
    _run(_para(tf, first=True, line=1.02), text, font=DISPLAY, size=size,
         bold=True, color=color, spacing=-0.9)
    return y + h + Inches(gap)


def lede(slide, text, y, *, size=17, width=None, color=INK_SOFT, italic=False,
         gap=0.26):
    w = width or Inches(9.9)
    h = block_height(text, size=size, width_in=w / Inches(1), line_spacing=1.38)
    tf = _text(slide, MARGIN, y, w, h + Inches(0.08))
    _run(_para(tf, first=True, line=1.35), text, font=BODY, size=size,
         color=color, italic=italic)
    return y + h + Inches(gap)


def bullets(slide, items, y, *, size=14, gap=13, width=None):
    """items: list of (bold_lead_or_None, rest)."""
    tf = _text(slide, MARGIN, y, width or Inches(11.1), H - y - Inches(1.1))
    for i, (lead, rest) in enumerate(items):
        p = _para(tf, first=(i == 0), space_after=gap, line=1.32)
        _run(p, "— ", font=BODY, size=size, color=ACCENT, bold=True)
        if lead:
            _run(p, lead, font=BODY, size=size, bold=True, color=INK)
            _run(p, "  " + rest, font=BODY, size=size, color=INK_SOFT)
        else:
            _run(p, rest, font=BODY, size=size, color=INK_SOFT)
    return y


def card_row(slide, cards, y, *, height=None, title_size=14, body_size=12,
             accent_top=False, numbered=False, fill=RAISED, edge=RULE):
    """cards: list of (title, body). Lays them out evenly across the column.

    ``height`` is measured from the tallest card's content unless given
    explicitly, so a long body can no longer spill past the card edge.
    """
    n = len(cards)
    gap = Inches(0.24)
    w = (CONTENT_W - gap * (n - 1)) / n
    pad = Inches(0.26)
    inner_in = (w - 2 * pad) / Inches(1)

    if height is None:
        tallest = Emu(0)
        for title, body in cards:
            used = Inches(0)
            if numbered:
                used += Inches(9 * 1.2 / 72) + Pt(7)
            if title:
                used += block_height(title, font=DISPLAY, size=title_size, bold=True,
                                     width_in=inner_in, line_spacing=1.22) + Pt(7)
            for line in (body if isinstance(body, list) else [body]):
                text = f"{line[0]}  {line[1]}" if isinstance(line, tuple) else line
                used += block_height(text, size=body_size, width_in=inner_in,
                                     line_spacing=1.34) + Pt(6)
            tallest = max(tallest, used)
        height = tallest + 2 * pad + Inches(0.06)
    for i, (title, body) in enumerate(cards):
        x = MARGIN + (w + gap) * i
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, height)
        box.adjustments[0] = 0.035
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = edge
        box.line.width = Pt(0.75)
        box.shadow.inherit = False
        box.text_frame.text = ""

        if accent_top:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.02), y,
                                         w - Inches(0.04), Inches(0.035))
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT
            bar.line.fill.background()
            bar.shadow.inherit = False

        tf = _text(slide, x + pad, y + pad, w - 2 * pad, height - 2 * pad)
        first = True
        if numbered:
            _run(_para(tf, first=True, space_after=7), f"{i + 1:02d}", font=MONO,
                 size=9, color=ACCENT, spacing=1.5)
            first = False
        if title:
            _run(_para(tf, first=first, space_after=7, line=1.2), title,
                 font=DISPLAY, size=title_size, bold=True, color=INK, spacing=-0.3)
            first = False
        for j, line in enumerate(body if isinstance(body, list) else [body]):
            p = _para(tf, first=first and j == 0, space_after=6, line=1.32)
            if isinstance(line, tuple):
                _run(p, line[0], font=BODY, size=body_size, bold=True, color=INK)
                _run(p, "  " + line[1], font=BODY, size=body_size, color=INK_SOFT)
            else:
                _run(p, line, font=BODY, size=body_size, color=INK_SOFT)
    return y + height + Inches(0.3)


def callout(slide, label, text, y, *, height=Inches(1.5)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, y, CONTENT_W, height)
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = CAUTION_WASH
    box.line.color.rgb = CAUTION
    box.line.width = Pt(0.75)
    box.shadow.inherit = False

    pad = Inches(0.3)
    tf = _text(slide, MARGIN + pad, y + pad, CONTENT_W - 2 * pad, height - 2 * pad)
    _run(_para(tf, first=True, space_after=8), label.upper(), font=MONO, size=9,
         color=CAUTION, spacing=1.5)
    _run(_para(tf, line=1.32), text, font=BODY, size=13, color=INK)
    return y + height + Inches(0.3)


def pull_quote(slide, text, y, *, size=21):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, Inches(0.035), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = _text(slide, MARGIN + Inches(0.32), y, Inches(9.4), Inches(1.3))
    _run(_para(tf, first=True, line=1.28), text, font=BODY, size=size,
         italic=True, color=INK)
    return y + Inches(1.4)


def metrics(slide, items, y):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y, CONTENT_W, Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False
    y = y + Inches(0.26)
    n = len(items)
    w = CONTENT_W / n
    for i, (num, label) in enumerate(items):
        x = MARGIN + w * i
        tf = _text(slide, x, y, w - Inches(0.2), Inches(1.1))
        _run(_para(tf, first=True, space_after=6, line=1.0), num, font=DISPLAY,
             size=34, bold=True, color=INK, spacing=-1.2)
        _run(_para(tf), label.upper(), font=MONO, size=8.5, color=INK_FAINT, spacing=1.4)
    return y + Inches(1.2)


# ---------------------------------------------------------------- the deck
def build() -> pathlib.Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 01 — title
    s = new_slide(prs, number=1, section="Agentic Graph Analytics")
    y = eyebrow(s, "Agentic Graph Analytics", y=Inches(2.35))
    y = headline(s, "Your data already knows\nwho the key players are.", y, size=48)
    y = lede(s, "Connect a database, say what you want to know in plain English, and get "
                "a reviewed, reproducible analysis back — without writing a line of code.",
             y + Inches(0.1), size=18, width=Inches(9.8))

    # 02 — the opportunity
    s = new_slide(prs, number=2, section="The opportunity")
    y = eyebrow(s, "The opportunity")
    y = headline(s, "The questions spreadsheets can’t answer", y)
    y = lede(s, "Every organisation can tell you totals. The questions that change "
                "decisions are about relationships.", y)
    card_row(s, [
        ("Fraud & financial crime",
         "“Which accounts sit at the centre of unusual money movement?”"),
        ("Advertising & identity",
         "“Which audiences secretly overlap, so we pay twice to reach the same people?”"),
        ("Clinical research",
         "“Which trial sites cluster around the same investigators — and what if one leaves?”"),
        ("Supply chain & risk",
         "“Which suppliers are we exposed to three hops away, through companies we’ve never heard of?”"),
    ], y + Inches(0.1), body_size=12)
    lede(s, "None of these are answerable by filtering rows. They’re about the shape of "
            "the connections between things.", y + Inches(2.45), size=15)

    # 03 — why now
    s = new_slide(prs, number=3, section="Why now")
    y = eyebrow(s, "Why now")
    y = headline(s, "The maths was never the obstacle", y)
    y = card_row(s, [
        ("PageRank", "The technique that made Google work. Tells you what is influential "
                     "in a network — not what is biggest, what is best-connected to other "
                     "well-connected things."),
        ("Community detection", "Finds the clusters nobody drew on the org chart — groups "
                                "that behave as a unit whether or not anyone intended them to."),
        ("Path & reachability", "Traces how exposure travels. Who is two steps from a "
                                "sanctioned party; what breaks if this node disappears."),
    ], y + Inches(0.1))
    pull_quote(s, "These have been production-ready for a decade. What took weeks was the "
                  "specialist work between having a database and having an answer.", y + Inches(0.15))

    # 04 — the five stages
    s = new_slide(prs, number=4, section="The product")
    y = eyebrow(s, "The product")
    y = headline(s, "Database to defensible answer, in five stages", y)
    y = card_row(s, [
        ("Connect", "Verify reachability, permissions and compute availability before you "
                    "invest any time."),
        ("Discover", "The schema is read from the database. You don’t describe your data "
                     "model — it finds it."),
        ("Ask", "Upload your requirements document, or let the Copilot interview you about "
                "your actual graph."),
        ("Curate", "Agents propose analyses. You approve, adjust or reject before anything runs."),
        ("Run & publish", "Live pipeline, results back in your database, reports with full "
                          "lineage."),
    ], y + Inches(0.2), accent_top=True, numbered=True,
        title_size=13, body_size=11)
    lede(s, "One console and a growing library of reusable templates — instead of a "
            "bespoke project per engagement.", y + Inches(0.3), size=15)

    # 05 — discover
    s = new_slide(prs, number=5, section="Discover")
    y = eyebrow(s, "Stage 02 — Discover", color=ACCENT)
    y = headline(s, "It reads the map, however messy the map is", y)
    y = lede(s, "Real databases don’t look like tutorials. Getting this stage right is what "
                "everything downstream depends on.", y)
    bullets(s, [
        ("Many logical types in one container.",
         "Plenty of databases cram dozens of entity types into a single collection, told "
         "apart by a type field. Naive tooling reports “one entity type: Entities” — which "
         "is useless. This doesn’t."),
        ("Several graphs sharing a database.",
         "A document corpus beside an extracted knowledge graph beside a straightforward "
         "HR dataset. Each is profiled separately and classified by what it is for."),
        ("Cross-graph links.",
         "Where those graphs actually join up is discovered by sampling real edges, not "
         "guessed from naming conventions."),
        ("Physical layout.",
         "How the data is spread across machines, read from the database itself — because "
         "that decides whether an analysis needs an expensive shuffle or can stay local."),
    ], y + Inches(0.15))

    # 06 — ask
    s = new_slide(prs, number=6, section="Ask")
    y = eyebrow(s, "Stage 03 — Ask", color=ACCENT)
    y = headline(s, "Two ways in, depending on where you’re starting", y)
    y = card_row(s, [
        ("Upload what you already have",
         "Requirements are usually already written down — a brief, a PDF, a Word document. "
         "Upload it and the structured objectives, requirements and constraints are "
         "extracted from the prose."),
        ("Or let the Copilot interview you",
         "Not a generic questionnaire. It has already seen your entity types, your "
         "relationship types, your volumes and whether your data is partitioned by "
         "customer — so it asks about your graph."),
    ], y + Inches(0.15), title_size=15)
    lede(s, "Either way you end up with a versioned requirements document you can review, "
            "edit, approve and revisit. Every analysis traces back to it.", y + Inches(0.1))

    # 07 — curate
    s = new_slide(prs, number=7, section="Curate")
    y = eyebrow(s, "Stage 04 — Curate", color=ACCENT)
    y = headline(s, "Agents propose. People approve.", y)
    y = lede(s, "Given your requirements and your schema, agents propose concrete analytical "
                "questions worth asking of this graph — each becoming a template with a "
                "chosen algorithm, its parameters, and its scope.", y)
    card_row(s, [
        ("Nothing runs unreviewed",
         "Templates arrive as DRAFT and stay there until a person approves them."),
        ("Approved means immutable",
         "Editing an approved template creates a new version rather than mutating the old "
         "one, so a published report always points at exactly what produced it."),
        ("Every knob is visible",
         "Disagree with the algorithm choice or the parameters? Change them. The agents do "
         "the tedious mapping; you keep the judgement."),
    ], y + Inches(0.2))

    # 08 — run & publish
    s = new_slide(prs, number=8, section="Run & publish")
    y = eyebrow(s, "Stage 05 — Run & publish", color=ACCENT)
    y = headline(s, "Watch it work, then defend the number", y)
    y = card_row(s, [
        ("Three ways to run", [
            ("Quick", "one prompt, one report, no setup"),
            ("Guided", "Copilot interview, then a focused analysis"),
            ("Detailed", "full requirements, many use cases, branches in parallel"),
        ]),
        ("Live pipeline",
         "Every step with its status, artifacts and warnings. Retry, pause and resume — a "
         "failure at step four doesn’t mean starting over."),
        ("Reports that stay true",
         "Rendered from database records, not exported as files. No more circulating a PDF "
         "that’s four months stale."),
    ], y + Inches(0.1))
    lede(s, "Because everything is a record with lineage, you can compare this quarter "
            "against last, trace any figure back to the template, requirement version and "
            "run that produced it, and re-run months later with a defensible account of "
            "what changed.", y, size=14)

    # 09 — governance
    s = new_slide(prs, number=9, section="Governance")
    y = eyebrow(s, "Past the demo")
    y = headline(s, "Built for how deployments actually look", y)
    y = callout(s, "The one that saves you from yourself",
                "Many production databases hold many customers’ data in one place, kept apart "
                "by a single field. Run an analysis blind across that and you get a beautiful "
                "result that mixes customers together — worse than no result. The product "
                "detects that layout from the database, names the field it found, and warns "
                "you before the run starts.",
                y + Inches(0.05))
    card_row(s, [
        ("Secrets stay out",
         "Credentials are held as references, never values. Exported bundles exclude them."),
        ("Everything is audited",
         "Create, update, approve, launch, publish, import, export, archive — recorded and "
         "attributable."),
        ("Retention that shows its work",
         "The cleanup sweep is a dry run by default: it names every record it would remove "
         "before anything is deleted. Approved requirements, published snapshots and the runs "
         "behind them are never removed at any age."),
    ], y)

    # 10 — verticals
    s = new_slide(prs, number=10, section="Verticals")
    y = eyebrow(s, "Time to value")
    y = headline(s, "Start from a vertical, not from zero", y)
    y = lede(s, "Five industry verticals ship with specialised analysis prompts and pattern "
                "detectors. Custom verticals can be generated per project, and a domain’s use "
                "cases and templates travel as a portable project bundle.", y)
    # Five cards share the content width, so the copy here is deliberately
    # terser than the HTML deck's — long bodies wrap badly at ~1.9in.
    y = card_row(s, [
        ("Ad-Tech & identity",
         "Audience overlap, identity stitching, influence ranking."),
        ("FinTech",
         "Exposure paths, counterparty networks, concentration risk."),
        ("Fraud intelligence",
         "Mule-account rings, device and IP clustering."),
        ("Social networks",
         "Community structure, influence, information flow."),
        ("Generic",
         "A sensible default for any other domain."),
    ], y + Inches(0.1))
    lede(s, "Bundles import as drafts, never as approved work — a starter pack is a starting "
            "point, not something that quietly starts running. Importing cannot execute code: "
            "a file that tries to smuggle executable content is rejected rather than run. The "
            "format is documented, so your team’s accumulated templates become the bundle you "
            "hand to the next project.", y, size=13.5)

    # 11 — who it's for
    s = new_slide(prs, number=11, section="Get started")
    y = eyebrow(s, "Who it’s for")
    y = headline(s, "Ask the question. Keep the receipt.", y)
    y = card_row(s, [
        ("Solutions & data teams",
         "Currently spinning up a bespoke project per engagement. Trade that for one console "
         "and a growing template library."),
        ("Domain analysts",
         "You don’t need to know what a damping factor is to ask “who are the key players "
         "here?” Review the proposal in business terms."),
        ("Anyone defending a number",
         "Versioned requirements, immutable approved templates, run lineage and a full audit "
         "trail. “Where did this come from?” is a link, not an excavation."),
    ], y + Inches(0.05))
    y = metrics(s, [
        ("5", "stages, one console"),
        ("3", "run modes"),
        ("3", "vertical starter bundles"),
        ("0", "lines of code to write"),
    ], y - Inches(0.05))
    lede(s, "Runs on ArangoDB — a database that stores data as a network of connected "
            "things, with an engine built to run these algorithms at scale.",
         y - Inches(0.15), size=12, color=INK_FAINT)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"wrote {path}")
